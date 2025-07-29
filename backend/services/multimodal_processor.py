"""
Multi-Modal Content Processor
Handles processing of various content types including text, images, audio, video,
and structured data with advanced extraction and analysis capabilities.
"""
import asyncio
import logging
import json
import base64
import io
import tempfile
import os
from datetime import datetime
from typing import Dict, List, Optional, Any, Union, Tuple
from dataclasses import dataclass
from enum import Enum
import uuid
import re
from pathlib import Path

import numpy as np
import pandas as pd
from PIL import Image, ImageEnhance, ImageFilter
import cv2
import pytesseract
import speech_recognition as sr
import librosa
import matplotlib.pyplot as plt
import seaborn as sns
from moviepy.editor import VideoFileClip
import fitz  # PyMuPDF
import docx
import openpyxl
from pptx import Presentation
import easyocr
from transformers import pipeline, BlipProcessor, BlipForConditionalGeneration
import torch

from sqlalchemy.orm import Session
from sqlalchemy import func, and_, desc, or_

from core.database import (
    get_db, Document, DocumentChunk, DocumentTag, AnalyticsEvent,
    User, UserProfile
)
from models.schemas import DocumentTagCreate, AnalyticsEventCreate

logger = logging.getLogger(__name__)

class ContentType(str, Enum):
    """Content type enumeration"""
    TEXT = "text"
    IMAGE = "image"
    AUDIO = "audio"
    VIDEO = "video"
    PDF = "pdf"
    DOCUMENT = "document"
    SPREADSHEET = "spreadsheet"
    PRESENTATION = "presentation"
    CODE = "code"
    STRUCTURED_DATA = "structured_data"
    WEB_CONTENT = "web_content"

class ProcessingQuality(str, Enum):
    """Processing quality levels"""
    FAST = "fast"
    BALANCED = "balanced"
    HIGH_QUALITY = "high_quality"
    MAXIMUM = "maximum"

class ExtractionMethod(str, Enum):
    """Content extraction methods"""
    OCR = "ocr"
    SPEECH_TO_TEXT = "speech_to_text"
    IMAGE_CAPTIONING = "image_captioning"
    VIDEO_ANALYSIS = "video_analysis"
    DOCUMENT_PARSING = "document_parsing"
    DATA_EXTRACTION = "data_extraction"
    CODE_ANALYSIS = "code_analysis"

@dataclass
class ProcessingResult:
    """Result of content processing"""
    content_type: ContentType
    extracted_text: str
    metadata: Dict[str, Any]
    confidence_score: float
    processing_time: float
    extraction_methods: List[ExtractionMethod]
    structured_data: Optional[Dict[str, Any]] = None
    embeddings: Optional[List[float]] = None
    thumbnails: Optional[List[str]] = None  # Base64 encoded images
    annotations: Optional[List[Dict[str, Any]]] = None

@dataclass
class ImageAnalysisResult:
    """Result of image analysis"""
    description: str
    objects_detected: List[Dict[str, Any]]
    text_extracted: str
    faces_detected: int
    scene_classification: str
    color_analysis: Dict[str, Any]
    quality_metrics: Dict[str, float]
    confidence_score: float

@dataclass
class AudioAnalysisResult:
    """Result of audio analysis"""
    transcription: str
    language_detected: str
    speaker_count: int
    emotion_analysis: Dict[str, float]
    audio_quality: Dict[str, float]
    duration: float
    confidence_score: float

@dataclass
class VideoAnalysisResult:
    """Result of video analysis"""
    transcription: str
    scene_descriptions: List[str]
    key_frames: List[str]  # Base64 encoded
    objects_timeline: List[Dict[str, Any]]
    audio_analysis: AudioAnalysisResult
    duration: float
    resolution: Tuple[int, int]
    confidence_score: float

class MultiModalProcessor:
    """Main multi-modal content processor"""
    
    def __init__(self, db: Session):
        self.db = db
        
        # Initialize AI models
        self._init_models()
        
        # Processing configurations
        self.processing_configs = {
            ProcessingQuality.FAST: {
                "image_max_size": (800, 600),
                "audio_sample_rate": 16000,
                "video_fps": 1,
                "ocr_confidence": 60
            },
            ProcessingQuality.BALANCED: {
                "image_max_size": (1200, 900),
                "audio_sample_rate": 22050,
                "video_fps": 2,
                "ocr_confidence": 70
            },
            ProcessingQuality.HIGH_QUALITY: {
                "image_max_size": (1920, 1080),
                "audio_sample_rate": 44100,
                "video_fps": 5,
                "ocr_confidence": 80
            },
            ProcessingQuality.MAXIMUM: {
                "image_max_size": (3840, 2160),
                "audio_sample_rate": 48000,
                "video_fps": 10,
                "ocr_confidence": 90
            }
        }
        
        # Supported file extensions
        self.supported_extensions = {
            ContentType.IMAGE: ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'],
            ContentType.AUDIO: ['.mp3', '.wav', '.flac', '.aac', '.ogg', '.m4a'],
            ContentType.VIDEO: ['.mp4', '.avi', '.mov', '.wmv', '.flv', '.webm', '.mkv'],
            ContentType.PDF: ['.pdf'],
            ContentType.DOCUMENT: ['.docx', '.doc', '.txt', '.rtf', '.odt'],
            ContentType.SPREADSHEET: ['.xlsx', '.xls', '.csv', '.ods'],
            ContentType.PRESENTATION: ['.pptx', '.ppt', '.odp'],
            ContentType.CODE: ['.py', '.js', '.html', '.css', '.java', '.cpp', '.c', '.php', '.rb', '.go', '.rs']
        }

    def _init_models(self):
        """Initialize AI models for processing"""
        try:
            # Image captioning model
            self.image_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
            self.image_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
            
            # OCR reader
            self.ocr_reader = easyocr.Reader(['en'])
            
            # Speech recognition
            self.speech_recognizer = sr.Recognizer()
            
            # Emotion analysis pipeline
            self.emotion_analyzer = pipeline("text-classification", 
                                           model="j-hartmann/emotion-english-distilroberta-base")
            
            # Object detection pipeline
            self.object_detector = pipeline("object-detection", 
                                          model="facebook/detr-resnet-50")
            
            logger.info("AI models initialized successfully")
            
        except Exception as e:
            logger.error(f"Error initializing AI models: {str(e)}")
            # Fallback to basic processing
            self.image_processor = None
            self.image_model = None
            self.ocr_reader = None
            self.emotion_analyzer = None
            self.object_detector = None

    async def process_content(
        self,
        file_path: str,
        content_type: Optional[ContentType] = None,
        quality: ProcessingQuality = ProcessingQuality.BALANCED,
        user_id: Optional[str] = None
    ) -> ProcessingResult:
        """Process content from file path"""
        try:
            start_time = datetime.utcnow()
            
            # Determine content type if not provided
            if not content_type:
                content_type = self._detect_content_type(file_path)
            
            # Get processing configuration
            config = self.processing_configs[quality]
            
            # Process based on content type
            if content_type == ContentType.IMAGE:
                result = await self._process_image(file_path, config)
            elif content_type == ContentType.AUDIO:
                result = await self._process_audio(file_path, config)
            elif content_type == ContentType.VIDEO:
                result = await self._process_video(file_path, config)
            elif content_type == ContentType.PDF:
                result = await self._process_pdf(file_path, config)
            elif content_type == ContentType.DOCUMENT:
                result = await self._process_document(file_path, config)
            elif content_type == ContentType.SPREADSHEET:
                result = await self._process_spreadsheet(file_path, config)
            elif content_type == ContentType.PRESENTATION:
                result = await self._process_presentation(file_path, config)
            elif content_type == ContentType.CODE:
                result = await self._process_code(file_path, config)
            else:
                result = await self._process_text(file_path, config)
            
            # Calculate processing time
            processing_time = (datetime.utcnow() - start_time).total_seconds()
            result.processing_time = processing_time
            
            # Track analytics
            if user_id:
                await self._track_processing_event(user_id, content_type, quality, processing_time)
            
            return result
            
        except Exception as e:
            logger.error(f"Error processing content: {str(e)}")
            raise

    async def process_content_batch(
        self,
        file_paths: List[str],
        quality: ProcessingQuality = ProcessingQuality.BALANCED,
        user_id: Optional[str] = None,
        max_concurrent: int = 5
    ) -> List[ProcessingResult]:
        """Process multiple files concurrently"""
        try:
            semaphore = asyncio.Semaphore(max_concurrent)
            
            async def process_single(file_path: str) -> ProcessingResult:
                async with semaphore:
                    return await self.process_content(file_path, quality=quality, user_id=user_id)
            
            tasks = [process_single(file_path) for file_path in file_paths]
            results = await asyncio.gather(*tasks, return_exceptions=True)
            
            # Filter out exceptions and log errors
            valid_results = []
            for i, result in enumerate(results):
                if isinstance(result, Exception):
                    logger.error(f"Error processing {file_paths[i]}: {str(result)}")
                else:
                    valid_results.append(result)
            
            return valid_results
            
        except Exception as e:
            logger.error(f"Error in batch processing: {str(e)}")
            raise

    def _detect_content_type(self, file_path: str) -> ContentType:
        """Detect content type from file extension"""
        file_ext = Path(file_path).suffix.lower()
        
        for content_type, extensions in self.supported_extensions.items():
            if file_ext in extensions:
                return content_type
        
        return ContentType.TEXT  # Default fallback

    async def _process_image(self, file_path: str, config: Dict[str, Any]) -> ProcessingResult:
        """Process image content"""
        try:
            # Load and preprocess image
            image = Image.open(file_path)
            original_size = image.size
            
            # Resize if needed
            max_size = config["image_max_size"]
            if image.size[0] > max_size[0] or image.size[1] > max_size[1]:
                image.thumbnail(max_size, Image.Resampling.LANCZOS)
            
            # Convert to RGB if needed
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Extract text using OCR
            extracted_text = ""
            ocr_confidence = 0.0
            
            if self.ocr_reader:
                try:
                    ocr_results = self.ocr_reader.readtext(np.array(image))
                    text_parts = []
                    confidences = []
                    
                    for (bbox, text, conf) in ocr_results:
                        if conf > config["ocr_confidence"] / 100:
                            text_parts.append(text)
                            confidences.append(conf)
                    
                    extracted_text = " ".join(text_parts)
                    ocr_confidence = np.mean(confidences) if confidences else 0.0
                    
                except Exception as e:
                    logger.warning(f"OCR processing failed: {str(e)}")
            
            # Generate image description
            description = ""
            caption_confidence = 0.0
            
            if self.image_processor and self.image_model:
                try:
                    inputs = self.image_processor(image, return_tensors="pt")
                    out = self.image_model.generate(**inputs, max_length=50)
                    description = self.image_processor.decode(out[0], skip_special_tokens=True)
                    caption_confidence = 0.8  # Placeholder confidence
                    
                except Exception as e:
                    logger.warning(f"Image captioning failed: {str(e)}")
            
            # Detect objects
            objects_detected = []
            if self.object_detector:
                try:
                    detections = self.object_detector(image)
                    for detection in detections:
                        if detection['score'] > 0.5:  # Confidence threshold
                            objects_detected.append({
                                'label': detection['label'],
                                'confidence': detection['score'],
                                'bbox': detection['box']
                            })
                except Exception as e:
                    logger.warning(f"Object detection failed: {str(e)}")
            
            # Analyze image properties
            color_analysis = self._analyze_image_colors(image)
            quality_metrics = self._calculate_image_quality(image)
            
            # Create thumbnail
            thumbnail = self._create_thumbnail(image)
            
            # Combine extracted text and description
            combined_text = f"{description} {extracted_text}".strip()
            
            # Calculate overall confidence
            overall_confidence = (ocr_confidence + caption_confidence) / 2
            
            # Create metadata
            metadata = {
                "original_size": original_size,
                "processed_size": image.size,
                "format": image.format,
                "mode": image.mode,
                "description": description,
                "objects_detected": objects_detected,
                "color_analysis": color_analysis,
                "quality_metrics": quality_metrics,
                "ocr_confidence": ocr_confidence,
                "caption_confidence": caption_confidence
            }
            
            return ProcessingResult(
                content_type=ContentType.IMAGE,
                extracted_text=combined_text,
                metadata=metadata,
                confidence_score=overall_confidence,
                processing_time=0.0,  # Will be set by caller
                extraction_methods=[ExtractionMethod.OCR, ExtractionMethod.IMAGE_CAPTIONING],
                thumbnails=[thumbnail] if thumbnail else None
            )
            
        except Exception as e:
            logger.error(f"Error processing image: {str(e)}")
            raise

    async def _process_audio(self, file_path: str, config: Dict[str, Any]) -> ProcessingResult:
        """Process audio content"""
        try:
            # Load audio file
            audio_data, sample_rate = librosa.load(file_path, sr=config["audio_sample_rate"])
            duration = len(audio_data) / sample_rate
            
            # Convert to temporary WAV file for speech recognition
            with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as temp_file:
                temp_path = temp_file.name
                librosa.output.write_wav(temp_path, audio_data, sample_rate)
            
            try:
                # Speech to text
                transcription = ""
                confidence = 0.0
                
                with sr.AudioFile(temp_path) as source:
                    audio = self.speech_recognizer.record(source)
                    try:
                        transcription = self.speech_recognizer.recognize_google(audio)
                        confidence = 0.8  # Google API doesn't return confidence
                    except sr.UnknownValueError:
                        transcription = ""
                        confidence = 0.0
                    except sr.RequestError as e:
                        logger.warning(f"Speech recognition error: {str(e)}")
                        transcription = ""
                        confidence = 0.0
                
                # Analyze audio features
                audio_features = self._analyze_audio_features(audio_data, sample_rate)
                
                # Emotion analysis on transcription
                emotion_scores = {}
                if transcription and self.emotion_analyzer:
                    try:
                        emotions = self.emotion_analyzer(transcription)
                        for emotion in emotions:
                            emotion_scores[emotion['label']] = emotion['score']
                    except Exception as e:
                        logger.warning(f"Emotion analysis failed: {str(e)}")
                
                # Create metadata
                metadata = {
                    "duration": duration,
                    "sample_rate": sample_rate,
                    "channels": 1,  # Librosa loads as mono by default
                    "audio_features": audio_features,
                    "emotion_analysis": emotion_scores,
                    "transcription_confidence": confidence
                }
                
                return ProcessingResult(
                    content_type=ContentType.AUDIO,
                    extracted_text=transcription,
                    metadata=metadata,
                    confidence_score=confidence,
                    processing_time=0.0,
                    extraction_methods=[ExtractionMethod.SPEECH_TO_TEXT]
                )
                
            finally:
                # Clean up temporary file
                os.unlink(temp_path)
                
        except Exception as e:
            logger.error(f"Error processing audio: {str(e)}")
            raise

    async def _process_video(self, file_path: str, config: Dict[str, Any]) -> ProcessingResult:
        """Process video content"""
        try:
            # Load video
            video = VideoFileClip(file_path)
            duration = video.duration
            fps = config["video_fps"]
            
            # Extract audio and process
            audio_result = None
            if video.audio:
                with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as temp_audio:
                    audio_path = temp_audio.name
                    video.audio.write_audiofile(audio_path, verbose=False, logger=None)
                    
                try:
                    audio_result = await self._process_audio(audio_path, config)
                finally:
                    os.unlink(audio_path)
            
            # Extract key frames
            key_frames = []
            scene_descriptions = []
            frame_times = np.linspace(0, duration, min(10, int(duration * fps)))
            
            for frame_time in frame_times:
                try:
                    frame = video.get_frame(frame_time)
                    frame_image = Image.fromarray(frame.astype('uint8'))
                    
                    # Generate description for frame
                    if self.image_processor and self.image_model:
                        try:
                            inputs = self.image_processor(frame_image, return_tensors="pt")
                            out = self.image_model.generate(**inputs, max_length=30)
                            description = self.image_processor.decode(out[0], skip_special_tokens=True)
                            scene_descriptions.append(f"At {frame_time:.1f}s: {description}")
                        except Exception as e:
                            logger.warning(f"Frame captioning failed: {str(e)}")
                    
                    # Create thumbnail
                    thumbnail = self._create_thumbnail(frame_image)
                    if thumbnail:
                        key_frames.append(thumbnail)
                        
                except Exception as e:
                    logger.warning(f"Error processing frame at {frame_time}s: {str(e)}")
            
            # Combine transcription and scene descriptions
            combined_text = ""
            if audio_result:
                combined_text += audio_result.extracted_text + " "
            combined_text += " ".join(scene_descriptions)
            
            # Calculate confidence
            confidence = 0.0
            if audio_result:
                confidence = audio_result.confidence_score
            
            # Create metadata
            metadata = {
                "duration": duration,
                "fps": video.fps,
                "resolution": (video.w, video.h),
                "has_audio": video.audio is not None,
                "scene_descriptions": scene_descriptions,
                "key_frame_count": len(key_frames),
                "audio_analysis": audio_result.metadata if audio_result else None
            }
            
            video.close()
            
            return ProcessingResult(
                content_type=ContentType.VIDEO,
                extracted_text=combined_text.strip(),
                metadata=metadata,
                confidence_score=confidence,
                processing_time=0.0,
                extraction_methods=[ExtractionMethod.VIDEO_ANALYSIS, ExtractionMethod.SPEECH_TO_TEXT],
                thumbnails=key_frames
            )
            
        except Exception as e:
            logger.error(f"Error processing video: {str(e)}")
            raise

    async def _process_pdf(self, file_path: str, config: Dict[str, Any]) -> ProcessingResult:
        """Process PDF content"""
        try:
            doc = fitz.open(file_path)
            extracted_text = ""
            images_text = ""
            page_count = len(doc)
            
            # Extract text and images from each page
            for page_num in range(page_count):
                page = doc.load_page(page_num)
                
                # Extract text
                page_text = page.get_text()
                extracted_text += f"\\n--- Page {page_num + 1} ---\\n{page_text}"
                
                # Extract images and process with OCR
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        
                        if pix.n - pix.alpha < 4:  # GRAY or RGB
                            img_data = pix.tobytes("png")
                            img_pil = Image.open(io.BytesIO(img_data))
                            
                            # OCR on image
                            if self.ocr_reader:
                                try:
                                    ocr_results = self.ocr_reader.readtext(np.array(img_pil))
                                    img_text = " ".join([text for (_, text, conf) in ocr_results 
                                                       if conf > config["ocr_confidence"] / 100])
                                    if img_text:
                                        images_text += f"\\n[Image {img_index + 1} on page {page_num + 1}]: {img_text}"
                                except Exception as e:
                                    logger.warning(f"OCR on PDF image failed: {str(e)}")
                        
                        pix = None
                        
                    except Exception as e:
                        logger.warning(f"Error processing PDF image: {str(e)}")
            
            doc.close()
            
            # Combine text from pages and images
            combined_text = extracted_text + images_text
            
            # Create metadata
            metadata = {
                "page_count": page_count,
                "has_images": len(images_text) > 0,
                "text_length": len(extracted_text),
                "images_text_length": len(images_text)
            }
            
            return ProcessingResult(
                content_type=ContentType.PDF,
                extracted_text=combined_text,
                metadata=metadata,
                confidence_score=0.9,  # High confidence for PDF text extraction
                processing_time=0.0,
                extraction_methods=[ExtractionMethod.DOCUMENT_PARSING, ExtractionMethod.OCR]
            )
            
        except Exception as e:
            logger.error(f"Error processing PDF: {str(e)}")
            raise

    async def _process_document(self, file_path: str, config: Dict[str, Any]) -> ProcessingResult:
        """Process document files (DOCX, DOC, TXT, etc.)"""
        try:
            file_ext = Path(file_path).suffix.lower()
            extracted_text = ""
            
            if file_ext == '.docx':
                doc = docx.Document(file_path)
                paragraphs = [paragraph.text for paragraph in doc.paragraphs]
                extracted_text = "\\n".join(paragraphs)
                
                # Extract table content
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join([cell.text for cell in row.cells])
                        extracted_text += f"\\n{row_text}"
                
            elif file_ext in ['.txt', '.rtf']:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    extracted_text = f.read()
            
            else:
                # Fallback to text extraction
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    extracted_text = f.read()
            
            # Analyze document structure
            lines = extracted_text.split('\\n')
            word_count = len(extracted_text.split())
            char_count = len(extracted_text)
            
            # Create metadata
            metadata = {
                "file_extension": file_ext,
                "line_count": len(lines),
                "word_count": word_count,
                "character_count": char_count,
                "has_tables": "table" in file_ext or "|" in extracted_text
            }
            
            return ProcessingResult(
                content_type=ContentType.DOCUMENT,
                extracted_text=extracted_text,
                metadata=metadata,
                confidence_score=0.95,  # High confidence for document parsing
                processing_time=0.0,
                extraction_methods=[ExtractionMethod.DOCUMENT_PARSING]
            )
            
        except Exception as e:
            logger.error(f"Error processing document: {str(e)}")
            raise

    async def _process_spreadsheet(self, file_path: str, config: Dict[str, Any]) -> ProcessingResult:
        """Process spreadsheet files"""
        try:
            file_ext = Path(file_path).suffix.lower()
            
            if file_ext == '.csv':
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path, sheet_name=None)  # Read all sheets
                
                # If multiple sheets, combine them
                if isinstance(df, dict):
                    combined_df = pd.DataFrame()
                    sheet_info = {}
                    
                    for sheet_name, sheet_df in df.items():
                        sheet_info[sheet_name] = {
                            "rows": len(sheet_df),
                            "columns": len(sheet_df.columns),
                            "column_names": list(sheet_df.columns)
                        }
                        
                        # Add sheet identifier
                        sheet_df['_sheet_name'] = sheet_name
                        combined_df = pd.concat([combined_df, sheet_df], ignore_index=True)
                    
                    df = combined_df
                else:
                    sheet_info = {"single_sheet": {
                        "rows": len(df),
                        "columns": len(df.columns),
                        "column_names": list(df.columns)
                    }}
            
            # Convert to text representation
            extracted_text = f"Spreadsheet Summary:\\n"
            extracted_text += f"Rows: {len(df)}, Columns: {len(df.columns)}\\n"
            extracted_text += f"Column Names: {', '.join(df.columns)}\\n\\n"
            
            # Add sample data
            extracted_text += "Sample Data (first 5 rows):\\n"
            extracted_text += df.head().to_string(index=False)
            
            # Add statistical summary for numeric columns
            numeric_cols = df.select_dtypes(include=[np.number]).columns
            if len(numeric_cols) > 0:
                extracted_text += "\\n\\nNumeric Column Statistics:\\n"
                extracted_text += df[numeric_cols].describe().to_string()
            
            # Analyze data types and patterns
            data_analysis = {
                "total_rows": len(df),
                "total_columns": len(df.columns),
                "column_types": df.dtypes.to_dict(),
                "numeric_columns": list(numeric_cols),
                "missing_values": df.isnull().sum().to_dict(),
                "unique_values": {col: df[col].nunique() for col in df.columns}
            }
            
            # Create structured data representation
            structured_data = {
                "headers": list(df.columns),
                "sample_rows": df.head(10).to_dict('records'),
                "statistics": data_analysis
            }
            
            # Create metadata
            metadata = {
                "file_extension": file_ext,
                "sheet_info": sheet_info if 'sheet_info' in locals() else None,
                "data_analysis": data_analysis
            }
            
            return ProcessingResult(
                content_type=ContentType.SPREADSHEET,
                extracted_text=extracted_text,
                metadata=metadata,
                confidence_score=0.95,
                processing_time=0.0,
                extraction_methods=[ExtractionMethod.DATA_EXTRACTION],
                structured_data=structured_data
            )
            
        except Exception as e:
            logger.error(f"Error processing spreadsheet: {str(e)}")
            raise

    async def _process_presentation(self, file_path: str, config: Dict[str, Any]) -> ProcessingResult:
        """Process presentation files"""
        try:
            prs = Presentation(file_path)
            extracted_text = ""
            slide_count = len(prs.slides)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"\\n--- Slide {slide_num} ---\\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\\n"
                
                extracted_text += slide_text
            
            # Create metadata
            metadata = {
                "slide_count": slide_count,
                "average_text_per_slide": len(extracted_text) / slide_count if slide_count > 0 else 0
            }
            
            return ProcessingResult(
                content_type=ContentType.PRESENTATION,
                extracted_text=extracted_text,
                metadata=metadata,
                confidence_score=0.9,
                processing_time=0.0,
                extraction_methods=[ExtractionMethod.DOCUMENT_PARSING]
            )
            
        except Exception as e:
            logger.error(f"Error processing presentation: {str(e)}")
            raise

    async def _process_code(self, file_path: str, config: Dict[str, Any]) -> ProcessingResult:
        """Process code files"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                code_content = f.read()
            
            file_ext = Path(file_path).suffix.lower()
            
            # Analyze code structure
            lines = code_content.split('\\n')
            total_lines = len(lines)
            code_lines = len([line for line in lines if line.strip() and not line.strip().startswith('#')])
            comment_lines = total_lines - code_lines
            
            # Extract functions/classes (basic pattern matching)
            functions = re.findall(r'def\\s+(\\w+)\\s*\\(', code_content)  # Python functions
            classes = re.findall(r'class\\s+(\\w+)\\s*[:\\(]', code_content)  # Python classes
            
            # Language-specific analysis
            language_features = self._analyze_code_language(code_content, file_ext)
            
            # Create enhanced text representation
            extracted_text = f"Code File Analysis:\\n"
            extracted_text += f"Language: {language_features['language']}\\n"
            extracted_text += f"Total Lines: {total_lines}\\n"
            extracted_text += f"Code Lines: {code_lines}\\n"
            extracted_text += f"Comment Lines: {comment_lines}\\n"
            
            if functions:
                extracted_text += f"Functions: {', '.join(functions)}\\n"
            if classes:
                extracted_text += f"Classes: {', '.join(classes)}\\n"
            
            extracted_text += f"\\nCode Content:\\n{code_content}"
            
            # Create metadata
            metadata = {
                "file_extension": file_ext,
                "language": language_features['language'],
                "total_lines": total_lines,
                "code_lines": code_lines,
                "comment_lines": comment_lines,
                "functions": functions,
                "classes": classes,
                "language_features": language_features
            }
            
            return ProcessingResult(
                content_type=ContentType.CODE,
                extracted_text=extracted_text,
                metadata=metadata,
                confidence_score=0.95,
                processing_time=0.0,
                extraction_methods=[ExtractionMethod.CODE_ANALYSIS]
            )
            
        except Exception as e:
            logger.error(f"Error processing code: {str(e)}")
            raise

    async def _process_text(self, file_path: str, config: Dict[str, Any]) -> ProcessingResult:
        """Process plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                extracted_text = f.read()
            
            # Basic text analysis
            lines = extracted_text.split('\\n')
            words = extracted_text.split()
            
            # Create metadata
            metadata = {
                "line_count": len(lines),
                "word_count": len(words),
                "character_count": len(extracted_text),
                "encoding": "utf-8"
            }
            
            return ProcessingResult(
                content_type=ContentType.TEXT,
                extracted_text=extracted_text,
                metadata=metadata,
                confidence_score=1.0,
                processing_time=0.0,
                extraction_methods=[ExtractionMethod.DOCUMENT_PARSING]
            )
            
        except Exception as e:
            logger.error(f"Error processing text: {str(e)}")
            raise

    # Helper methods
    def _analyze_image_colors(self, image: Image.Image) -> Dict[str, Any]:
        """Analyze color properties of an image"""
        try:
            # Convert to RGB if needed
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # Get dominant colors
            colors = image.getcolors(maxcolors=256*256*256)
            if colors:
                dominant_color = max(colors, key=lambda x: x[0])[1]
                
                # Calculate color statistics
                pixels = np.array(image)
                mean_color = np.mean(pixels, axis=(0, 1))
                std_color = np.std(pixels, axis=(0, 1))
                
                return {
                    "dominant_color": dominant_color,
                    "mean_color": mean_color.tolist(),
                    "std_color": std_color.tolist(),
                    "color_count": len(colors)
                }
            
            return {"error": "Could not analyze colors"}
            
        except Exception as e:
            logger.warning(f"Color analysis failed: {str(e)}")
            return {"error": str(e)}

    def _calculate_image_quality(self, image: Image.Image) -> Dict[str, float]:
        """Calculate image quality metrics"""
        try:
            # Convert to grayscale for analysis
            gray = image.convert('L')
            pixels = np.array(gray)
            
            # Calculate sharpness (Laplacian variance)
            laplacian = cv2.Laplacian(pixels, cv2.CV_64F)
            sharpness = laplacian.var()
            
            # Calculate brightness
            brightness = np.mean(pixels)
            
            # Calculate contrast
            contrast = np.std(pixels)
            
            return {
                "sharpness": float(sharpness),
                "brightness": float(brightness),
                "contrast": float(contrast)
            }
            
        except Exception as e:
            logger.warning(f"Quality calculation failed: {str(e)}")
            return {"error": str(e)}

    def _create_thumbnail(self, image: Image.Image, size: Tuple[int, int] = (150, 150)) -> Optional[str]:
        """Create base64 encoded thumbnail"""
        try:
            thumbnail = image.copy()
            thumbnail.thumbnail(size, Image.Resampling.LANCZOS)
            
            buffer = io.BytesIO()
            thumbnail.save(buffer, format='JPEG', quality=85)
            buffer.seek(0)
            
            thumbnail_b64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
            return f"data:image/jpeg;base64,{thumbnail_b64}"
            
        except Exception as e:
            logger.warning(f"Thumbnail creation failed: {str(e)}")
            return None

    def _analyze_audio_features(self, audio_data: np.ndarray, sample_rate: int) -> Dict[str, Any]:
        """Analyze audio features"""
        try:
            # Extract features
            features = {}
            
            # Spectral features
            spectral_centroids = librosa.feature.spectral_centroid(y=audio_data, sr=sample_rate)[0]
            features['spectral_centroid_mean'] = float(np.mean(spectral_centroids))
            
            # MFCC features
            mfccs = librosa.feature.mfcc(y=audio_data, sr=sample_rate, n_mfcc=13)
            features['mfcc_mean'] = np.mean(mfccs, axis=1).tolist()
            
            # Zero crossing rate
            zcr = librosa.feature.zero_crossing_rate(audio_data)[0]
            features['zero_crossing_rate_mean'] = float(np.mean(zcr))
            
            # RMS energy
            rms = librosa.feature.rms(y=audio_data)[0]
            features['rms_mean'] = float(np.mean(rms))
            
            return features
            
        except Exception as e:
            logger.warning(f"Audio feature analysis failed: {str(e)}")
            return {"error": str(e)}

    def _analyze_code_language(self, code_content: str, file_ext: str) -> Dict[str, Any]:
        """Analyze code language and features"""
        language_map = {
            '.py': 'Python',
            '.js': 'JavaScript',
            '.html': 'HTML',
            '.css': 'CSS',
            '.java': 'Java',
            '.cpp': 'C++',
            '.c': 'C',
            '.php': 'PHP',
            '.rb': 'Ruby',
            '.go': 'Go',
            '.rs': 'Rust'
        }
        
        language = language_map.get(file_ext, 'Unknown')
        
        # Basic complexity analysis
        complexity_indicators = {
            'nested_blocks': len(re.findall(r'\\s{4,}if|\\s{4,}for|\\s{4,}while', code_content)),
            'function_calls': len(re.findall(r'\\w+\\s*\\(', code_content)),
            'imports': len(re.findall(r'import\\s+|from\\s+.*import', code_content)),
            'comments': len(re.findall(r'#.*|//.*|/\\*.*\\*/', code_content))
        }
        
        return {
            'language': language,
            'complexity_indicators': complexity_indicators
        }

    async def _track_processing_event(
        self, user_id: str, content_type: ContentType, 
        quality: ProcessingQuality, processing_time: float
    ):
        """Track processing analytics event"""
        try:
            event = AnalyticsEvent(
                user_id=user_id,
                event_type="multimodal_processing",
                event_data={
                    "content_type": content_type.value,
                    "quality": quality.value,
                    "processing_time": processing_time,
                    "timestamp": datetime.utcnow().isoformat(),
                    "service": "multimodal_processor"
                }
            )
            
            self.db.add(event)
            self.db.commit()
            
        except Exception as e:
            logger.error(f"Error tracking processing event: {str(e)}")

# Export classes
__all__ = [
    'MultiModalProcessor',
    'ProcessingResult',
    'ImageAnalysisResult',
    'AudioAnalysisResult',
    'VideoAnalysisResult',
    'ContentType',
    'ProcessingQuality',
    'ExtractionMethod'
]